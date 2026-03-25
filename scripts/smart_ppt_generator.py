#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
智能 PPT 生成器 - 根据主题自动设计配色和布局
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import os

# ==================== 主题配色方案 ====================

THEMES = {
    'tech': {
        'name': '科技蓝',
        'primary': RGBColor(25, 118, 210),
        'secondary': RGBColor(0, 188, 212),
        'accent': RGBColor(255, 152, 0),
        'success': RGBColor(76, 175, 80),
        'bg_light': RGBColor(245, 248, 255),
        'bg_dark': RGBColor(13, 71, 161),
        'text_dark': RGBColor(33, 33, 33),
        'text_light': RGBColor(255, 255, 255),
    },
    'business': {
        'name': '商务深蓝',
        'primary': RGBColor(23, 105, 170),
        'secondary': RGBColor(84, 110, 122),
        'accent': RGBColor(214, 112, 21),
        'success': RGBColor(102, 187, 106),
        'bg_light': RGBColor(245, 247, 250),
        'bg_dark': RGBColor(20, 50, 90),
        'text_dark': RGBColor(33, 33, 33),
        'text_light': RGBColor(255, 255, 255),
    },
    'creative': {
        'name': '创意活力',
        'primary': RGBColor(156, 39, 176),
        'secondary': RGBColor(33, 150, 243),
        'accent': RGBColor(255, 87, 34),
        'success': RGBColor(76, 175, 80),
        'bg_light': RGBColor(250, 245, 255),
        'bg_dark': RGBColor(74, 20, 140),
        'text_dark': RGBColor(33, 33, 33),
        'text_light': RGBColor(255, 255, 255),
    },
    'minimal': {
        'name': '极简黑白',
        'primary': RGBColor(0, 0, 0),
        'secondary': RGBColor(100, 100, 100),
        'accent': RGBColor(255, 87, 34),
        'success': RGBColor(76, 175, 80),
        'bg_light': RGBColor(255, 255, 255),
        'bg_dark': RGBColor(33, 33, 33),
        'text_dark': RGBColor(33, 33, 33),
        'text_light': RGBColor(255, 255, 255),
    },
}

# Emoji 图标库
EMOJI_ICONS = {
    'tech': '💻', 'ai': '🤖', 'data': '📊', 'security': '🔒',
    'cloud': '☁️', 'mobile': '📱', 'network': '🌐', 'code': '💻',
    'business': '💼', 'money': '💰', 'chart': '📈', 'team': '👥',
    'time': '⏰', 'location': '📍', 'email': '📧', 'phone': '📞',
    'check': '✅', 'warning': '⚠️', 'error': '❌', 'info': 'ℹ️',
    'star': '⭐', 'heart': '❤️', 'fire': '🔥', 'rocket': '🚀',
}

def get_theme_for_content(content_type):
    """根据内容类型自动选择主题"""
    theme_map = {
        'tech': ['技术', 'AI', '智能', '科技', '数字化', 'OpenClaw', 'Agent'],
        'business': ['商务', '企业', '应用', '方案', '商业', '价值'],
        'creative': ['创意', '设计', '营销', '品牌'],
        'minimal': ['简约', '报告', '总结'],
    }
    
    for theme_name, keywords in theme_map.items():
        for keyword in keywords:
            if keyword in content_type:
                return THEMES[theme_name]
    
    return THEMES['tech']  # 默认科技蓝

def add_title_slide(prs, title, subtitle, theme):
    """添加标题页 - 渐变背景设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加纯色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['primary']
    bg.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = theme['text_light']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, section_num, theme):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色块
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['bg_dark']
    bg.line.fill.background()
    
    # 章节编号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(2))
    tf = num_box.text_frame
    p = tf.add_paragraph()
    p.text = str(section_num)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = theme['accent']
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(9), Inches(3))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    
    return slide

def add_content_slide(prs, title, content_items, theme, show_emoji=True):
    """添加内容页 - 卡片式布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    
    # 内容区域 - 卡片式布局
    start_y = Inches(1.5)
    card_height = Inches(1.2)
    spacing = Inches(0.3)
    
    for i, item in enumerate(content_items):
        card_y = start_y + i * (card_height + spacing)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), card_y, Inches(12.333), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = theme['bg_light']
        card.line.color.rgb = theme['primary']
        card.line.width = Pt(1)
        
        # 添加 Emoji 图标
        if show_emoji and item.get('emoji'):
            emoji_box = slide.shapes.add_textbox(
                Inches(0.7), card_y + Inches(0.2), Inches(0.5), Inches(0.8)
            )
            tf = emoji_box.text_frame
            p = tf.add_paragraph()
            p.text = item['emoji']
            p.font.size = Pt(28)
        
        # 添加内容
        text_x = Inches(1.3) if show_emoji and item.get('emoji') else Inches(0.7)
        content_box = slide.shapes.add_textbox(
            text_x, card_y + Inches(0.2), Inches(11), Inches(0.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = item['text']
        p.font.size = Pt(18)
        p.font.color.rgb = theme['text_dark']
        
        if item.get('bold'):
            p.font.bold = True
    
    return slide

def add_comparison_table(prs, title, columns, rows, theme):
    """添加对比表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(11), Inches(0.6))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    
    # 创建表格
    table_rows = len(rows) + 1
    table_cols = len(columns)
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.333)
    height = Inches(0.6)
    
    table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
    
    # 设置列宽
    col_width = int(width / table_cols)
    for i in range(table_cols):
        table.columns[i].width = col_width
    
    # 填充表头
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme['primary']
        cell.text_frame.paragraphs[0].font.color.rgb = theme['text_light']
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # 填充内容
    for row_idx, row in enumerate(rows, 1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            
            # 隔行换色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme['bg_light']
            
            # 第一列加粗
            if col_idx == 0:
                cell.text_frame.paragraphs[0].font.bold = True
    
    return slide

def add_chart_slide(prs, title, chart_data, chart_type='bar', theme=None):
    """添加图表页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary'] if theme else THEMES['tech']['primary']
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(11), Inches(0.6))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme['text_light'] if theme else THEMES['tech']['text_light']
    
    # 创建图表
    fig, ax = plt.subplots(figsize=(10, 5))
    
    if chart_type == 'bar':
        ax.bar(chart_data['labels'], chart_data['values'], color='#1976D2')
    elif chart_type == 'line':
        ax.plot(chart_data['labels'], chart_data['values'], color='#1976D2', marker='o')
    elif chart_type == 'pie':
        ax.pie(chart_data['values'], labels=chart_data['labels'], autopct='%1.1f%%')
    
    ax.set_xlabel('Category')
    ax.set_ylabel('Value')
    ax.tick_params(axis='x', rotation=45)
    plt.tight_layout()
    
    # 保存图表
    chart_path = '/tmp/chart_temp.png'
    plt.savefig(chart_path, dpi=150, bbox_inches='tight')
    plt.close()
    
    # 插入到 PPT
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(12.333))
    
    # 清理临时文件
    try:
        os.remove(chart_path)
    except:
        pass
    
    return slide

def add_summary_slide(prs, title, key_points, theme):
    """添加总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['bg_light']
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = theme['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # 关键点
    start_y = Inches(2)
    for i, point in enumerate(key_points):
        y_pos = start_y + i * Inches(1)
        
        # 添加序号
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), y_pos, Inches(0.8), Inches(0.8)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = theme['primary']
        num_circle.line.fill.background()
        
        num_text = slide.shapes.add_textbox(Inches(1.15), y_pos + Inches(0.2), Inches(0.5), Inches(0.4))
        tf = num_text.text_frame
        p = tf.add_paragraph()
        p.text = str(i + 1)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = theme['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 添加内容
        content_box = slide.shapes.add_textbox(Inches(2), y_pos + Inches(0.2), Inches(10), Inches(0.6))
        tf = content_box.text_frame
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = theme['text_dark']
    
    return slide

# ==================== 主函数 ====================

def create_smart_ppt(title, subtitle, slides_content, output_path, content_type='tech'):
    """智能创建 PPT"""
    print(f"创建智能 PPT: {title}")
    
    # 自动选择主题
    theme = get_theme_for_content(content_type)
    print(f"使用主题：{theme['name']}")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    section_num = 1
    
    for slide_data in slides_content:
        slide_type = slide_data.get('type', 'content')
        
        if slide_type == 'title':
            print(f"  创建标题页...")
            add_title_slide(prs, slide_data['title'], slide_data.get('subtitle', subtitle), theme)
        
        elif slide_type == 'section':
            print(f"  创建章节页：{slide_data['title']}...")
            add_section_slide(prs, slide_data['title'], section_num, theme)
            section_num += 1
        
        elif slide_type == 'content':
            print(f"  创建内容页：{slide_data['title']}...")
            add_content_slide(prs, slide_data['title'], slide_data['items'], theme)
        
        elif slide_type == 'comparison':
            print(f"  创建对比页：{slide_data['title']}...")
            add_comparison_table(prs, slide_data['title'], slide_data['columns'], slide_data['rows'], theme)
        
        elif slide_type == 'chart':
            print(f"  创建图表页：{slide_data['title']}...")
            add_chart_slide(prs, slide_data['title'], slide_data['data'], slide_data.get('chart_type', 'bar'), theme)
        
        elif slide_type == 'summary':
            print(f"  创建总结页：{slide_data['title']}...")
            add_summary_slide(prs, slide_data['title'], slide_data['points'], theme)
    
    # 保存
    prs.save(output_path)
    print(f"\n✅ PPT 已保存：{output_path}")
    print(f"📊 共 {len(slides_content)} 页幻灯片")
    
    return output_path

if __name__ == '__main__':
    # 测试数据
    slides_content = [
        {'type': 'title', 'title': '测试 PPT', 'subtitle': '智能生成'},
        {'type': 'section', 'title': '第一部分'},
        {'type': 'content', 'title': '内容页', 'items': [
            {'text': '要点 1', 'emoji': '✅'},
            {'text': '要点 2', 'emoji': '⭐'},
        ]},
    ]
    
    create_smart_ppt(
        '测试 PPT',
        '智能生成',
        slides_content,
        '/Users/yangbowen/Desktop/测试 PPT.pptx',
        'tech'
    )
