#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
专业级 PPT 生成器 - 遵循设计四原则
对齐 · 对比 · 重复 · 亲密性
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import os

# ==================== 专业配色方案 ====================

PROFESSIONAL_THEMES = {
    'china_mobile': {
        'name': '中国移动蓝',
        'primary': RGBColor(0, 104, 180),      # 移动蓝
        'secondary': RGBColor(0, 153, 204),    # 科技蓝
        'accent': RGBColor(255, 102, 0),       # 活力橙
        'success': RGBColor(0, 153, 102),      # 成功绿
        'bg_light': RGBColor(245, 248, 250),   # 浅灰背景
        'bg_dark': RGBColor(0, 51, 102),       # 深蓝背景
        'text_dark': RGBColor(51, 51, 51),     # 深灰文字
        'text_light': RGBColor(255, 255, 255), # 白色文字
        'border': RGBColor(200, 200, 200),     # 边框灰
    },
}

# ==================== 布局常量 ====================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
HEADER_HEIGHT = Inches(1.0)
FOOTER_HEIGHT = Inches(0.5)

# 字体大小
FONT_TITLE = Pt(36)
FONT_HEADING = Pt(28)
FONT_SUBHEADING = Pt(22)
FONT_BODY = Pt(18)
FONT_CAPTION = Pt(14)

# 行间距
LINE_SPACING = 1.3

def add_header(slide, title, theme):
    """添加统一的页眉"""
    # 顶部色条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, HEADER_HEIGHT
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme['primary']
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(0.25), CONTENT_WIDTH, Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = FONT_HEADING
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    
    return header

def add_footer(slide, page_num, total_pages, theme):
    """添加统一的页脚"""
    # 底部色条
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_HEIGHT - FOOTER_HEIGHT, SLIDE_WIDTH, FOOTER_HEIGHT
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = theme['bg_dark']
    footer.line.fill.background()
    
    # 页码
    page_text = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.4), Inches(1), Inches(0.3)
    )
    tf = page_text.text_frame
    p = tf.add_paragraph()
    p.text = f"{page_num}/{total_pages}"
    p.font.size = FONT_CAPTION
    p.font.color.rgb = theme['text_light']
    p.alignment = PP_ALIGN.RIGHT
    
    return footer

def add_title_slide(prs, title, subtitle, theme):
    """标题页 - 渐变背景 + 居中布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 纯色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['primary']
    bg.line.fill.background()
    
    # 标题（居中）
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), CONTENT_WIDTH, Inches(2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.size = FONT_TITLE
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题（居中）
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5), CONTENT_WIDTH, Inches(1)
    )
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = FONT_SUBHEADING
    p.font.color.rgb = theme['text_light']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, section_title, theme):
    """章节页 - 深色背景 + 大字号"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 深色背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme['bg_dark']
    bg.line.fill.background()
    
    # 章节编号（超大字号）
    num_box = slide.shapes.add_textbox(
        MARGIN, Inches(1), Inches(3), Inches(3)
    )
    tf = num_box.text_frame
    p = tf.add_paragraph()
    p.text = str(section_num)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = theme['accent']
    
    # 章节标题
    title_box = slide.shapes.add_textbox(
        Inches(4), Inches(2), CONTENT_WIDTH - Inches(3), Inches(2)
    )
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = section_title
    p.font.size = FONT_TITLE
    p.font.bold = True
    p.font.color.rgb = theme['text_light']
    
    return slide

def add_content_slide(prs, title, items, theme, page_num, total_pages):
    """内容页 - 卡片式布局 + 严格对齐"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加页眉
    add_header(slide, title, theme)
    
    # 内容区域
    start_y = HEADER_HEIGHT + Inches(0.5)
    card_height = Inches(1.0)
    card_spacing = Inches(0.3)
    
    for i, item in enumerate(items):
        card_y = start_y + i * (card_height + card_spacing)
        
        # 卡片背景（圆角矩形）
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN, card_y, CONTENT_WIDTH, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = theme['bg_light']
        card.line.color.rgb = theme['border']
        card.line.width = Pt(1)
        
        # 内容文本框
        content_box = slide.shapes.add_textbox(
            MARGIN + Inches(0.3), card_y + Inches(0.2), CONTENT_WIDTH - Inches(0.6), Inches(0.6)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = item.get('text', '')
        p.font.size = FONT_BODY
        p.font.color.rgb = theme['text_dark']
        
        if item.get('bold'):
            p.font.bold = True
    
    # 添加页脚
    add_footer(slide, page_num, total_pages, theme)
    
    return slide

def add_comparison_slide(prs, title, columns, rows, theme, page_num, total_pages):
    """对比页 - 专业表格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加页眉
    add_header(slide, title, theme)
    
    # 创建表格
    table_rows = len(rows) + 1
    table_cols = len(columns)
    left = MARGIN
    top = HEADER_HEIGHT + Inches(0.5)
    width = CONTENT_WIDTH
    row_height = Inches(0.5)
    
    table = slide.shapes.add_table(
        table_rows, table_cols, left, top, width, row_height
    ).table
    
    # 设置列宽（平均分配）
    col_width = int(width / table_cols)
    for i in range(table_cols):
        table.columns[i].width = col_width
    
    # 填充表头
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme['primary']
        cell.text_frame.paragraphs[0].font.color.rgb = theme['text_light']
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = FONT_BODY
    
    # 填充内容
    for row_idx, row in enumerate(rows, 1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = FONT_BODY
            
            # 第一列加粗
            if col_idx == 0:
                cell.text_frame.paragraphs[0].font.bold = True
            
            # 隔行换色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme['bg_light']
    
    # 添加页脚
    add_footer(slide, page_num, total_pages, theme)
    
    return slide

def add_summary_slide(prs, title, key_points, theme):
    """总结页 - 编号列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加页眉
    add_header(slide, title, theme)
    
    # 关键点列表
    start_y = HEADER_HEIGHT + Inches(0.5)
    item_height = Inches(1.2)
    item_spacing = Inches(0.3)
    
    for i, point in enumerate(key_points):
        y_pos = start_y + i * (item_height + item_spacing)
        
        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            MARGIN, y_pos, Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme['primary']
        circle.line.fill.background()
        
        # 编号文字
        num_text = slide.shapes.add_textbox(
            MARGIN + Inches(0.15), y_pos + Inches(0.2), Inches(0.5), Inches(0.4)
        )
        tf = num_text.text_frame
        p = tf.add_paragraph()
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = theme['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 内容文字
        content_box = slide.shapes.add_textbox(
            MARGIN + Inches(1), y_pos + Inches(0.3), CONTENT_WIDTH - Inches(1), Inches(0.6)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = point
        p.font.size = FONT_BODY
        p.font.color.rgb = theme['text_dark']
    
    return slide

# ==================== 主函数 ====================

def create_professional_ppt(output_path):
    """创建专业级 PPT"""
    print("创建专业级 OpenClaw 企业应用方案 PPT...")
    
    theme = PROFESSIONAL_THEMES['china_mobile']
    print(f"使用主题：{theme['name']}")
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 定义所有内容页
    slides_content = [
        # 1. 标题页
        ('title', 'OpenClaw 企业应用方案', '技术架构 · 应用场景 · 部署方案\n中国移动终端公司 · 2026 年 3 月'),
        
        # 2. 目录
        ('content', '目录', [
            {'text': 'OpenClaw 技术架构与原理'},
            {'text': 'OpenClaw vs 对话式 AI vs 普通 Agent vs RPA'},
            {'text': '中国移动终端公司应用场景'},
            {'text': 'OpenClaw + RPA 协同方案'},
            {'text': '为一线创造的价值'},
            {'text': '内网部署硬件要求（70B+ 大模型）'},
            {'text': '内网环境限制与风险'},
            {'text': '总结与建议'},
        ]),
        
        # 3. 章节页
        ('section', 1, '技术架构与原理'),
        
        # 4. 技术架构
        ('content', 'OpenClaw 技术架构', [
            {'text': '用户交互层：飞书/微信/Telegram/网页 - 多渠道接入'},
            {'text': 'Gateway 网关层：消息路由·会话管理·安全认证·限流控制'},
            {'text': 'Agent 核心层：意图识别·任务规划·工具调用·记忆检索·结果生成'},
            {'text': '技能工具层：文件操作·浏览器控制·API 调用·代码执行·数据库查询'},
            {'text': '记忆系统：短期记忆 (Session)·长期记忆 (MEMORY.md)·向量检索'},
        ]),
        
        # 5. 技术原理
        ('content', '核心设计理念', [
            {'text': '本地优先（Local-First）：所有数据存储在本地，保护隐私安全'},
            {'text': '极简内核 + 弹性扩展：核心代码精简，通过 Skills 扩展能力'},
            {'text': '多模型兼容：支持百炼、Gemini、GPT、Claude 等主流模型'},
            {'text': '', 'bold': False},
            {'text': 'Agent 循环（Agent Loop）：感知 → 思考 → 行动 → 反馈 → 学习', 'bold': True},
        ]),
        
        # 6. 章节页
        ('section', 2, '技术对比'),
        
        # 7. 对比表格
        ('comparison', 'OpenClaw vs 其他技术', 
         ['对比维度', '对话式 AI', '普通 Agent', 'RPA', 'OpenClaw'],
         [
             ['核心能力', '文本对话', '简单任务', '规则自动化', '系统级任务'],
             ['上下文理解', '短期对话', '有限', '无', '长期记忆'],
             ['工具调用', '❌', '部分', '固定脚本', '✅ 丰富 Skills'],
             ['学习能力', '❌', '弱', '❌', '✅ 自我改进'],
             ['主动性', '被动响应', '被动', '定时触发', '主动代理'],
             ['数据隐私', '低', '中', '高', '高'],
         ]),
        
        # 8. 章节页
        ('section', 3, '应用场景'),
        
        # 9. 业务部门
        ('content', '业务部门应用场景', [
            {'text': '市场部：竞品分析·舆情监控·报告生成 - 节省 87% 时间'},
            {'text': '产品部：需求分析·竞品对比·文档管理 - 提升产品决策效率'},
            {'text': '客服部：智能问答·工单处理·满意度分析 - 83% 效率提升'},
            {'text': '财务部：发票处理·报表生成·对账核对 - 节省 85% 时间'},
        ]),
        
        # 10. 支撑部门
        ('content', '支撑部门应用场景', [
            {'text': '供应链：库存监控·供应商管理·采购分析 - 92% 效率提升'},
            {'text': '人力部：简历筛选·培训管理·员工问答 - 自动化 HR 流程'},
            {'text': 'IT 运维：监控告警·日志分析·自动化运维 - 7x24 小时监控'},
            {'text': '研发部：代码审查·文档生成·Bug 分析 - 提升开发效率'},
        ]),
        
        # 11. 章节页
        ('section', 4, '价值评估'),
        
        # 12. 量化价值
        ('comparison', '量化价值评估',
         ['场景', '当前耗时', 'OpenClaw 后', '节省', '年化价值'],
         [
             ['销售日报', '2 小时/天', '10 分钟/天', '95%', '5 万元'],
             ['竞品分析', '8 小时/周', '1 小时/周', '87%', '15 万元'],
             ['发票录入', '4 小时/天', '30 分钟/天', '85%', '8 万元'],
             ['客服问答', '6 小时/天', '1 小时/天', '83%', '12 万元'],
             ['库存监控', '2 小时/天', '10 分钟/天', '92%', '4 万元'],
             ['合计', '-', '-', '-', '44 万元/年'],
         ]),
        
        # 13. 章节页
        ('section', 5, '部署方案'),
        
        # 14. 硬件要求
        ('content', '硬件配置要求（70B+ 大模型）', [
            {'text': 'GPU：4x A800 80GB（入门 80 万）或 8x H800 80GB（推荐 200 万）'},
            {'text': 'CPU：2x AMD EPYC 7763 或 4x EPYC 9654 - 64 核心+'},
            {'text': '内存：512GB-1TB DDR5 - 支持大模型推理'},
            {'text': '存储：2TB NVMe SSD + 10TB+ 数据盘 - 高速读写'},
            {'text': '', 'bold': False},
            {'text': '替代方案：16x RTX 4090 24GB（成本 50 万，性价比高）', 'bold': True},
        ]),
        
        # 15. 风险与缓解
        ('content', '内网限制与风险缓解', [
            {'text': '网络限制：无法访问外网 API → 使用本地模型或代理服务器'},
            {'text': '安全风险：数据泄露 → 沙箱隔离 + 权限控制 + 操作审计'},
            {'text': '合规风险：数据合规 → 数据脱敏 + 本地部署 + 人工审核'},
            {'text': '运维风险：系统依赖 → 冗余部署 + 监控告警 + 应急预案'},
        ]),
        
        # 16. 总结页
        ('summary', '总结与建议', [
            '核心优势：本地优先、系统级执行、丰富生态、灵活扩展',
            '实施建议：试点先行、场景优先、人机协同、安全先行',
            '下一步：需求调研 → 方案设计 → 环境准备 → 试点部署 → 效果评估',
            '推荐试点：市场部（竞品分析）、客服部（智能问答）',
        ]),
    ]
    
    # 计算总页数
    total_pages = len(slides_content)
    page_num = 1
    
    for slide_data in slides_content:
        slide_type = slide_data[0]
        print(f"  创建第{page_num}页：{slide_type}...")
        
        if slide_type == 'title':
            add_title_slide(prs, slide_data[1], slide_data[2], theme)
        
        elif slide_type == 'section':
            add_section_slide(prs, slide_data[1], slide_data[2], theme)
        
        elif slide_type == 'content':
            add_content_slide(prs, slide_data[1], slide_data[2], theme, page_num, total_pages)
        
        elif slide_type == 'comparison':
            add_comparison_slide(prs, slide_data[1], slide_data[2], slide_data[3], theme, page_num, total_pages)
        
        elif slide_type == 'summary':
            add_summary_slide(prs, slide_data[1], slide_data[2], theme)
        
        page_num += 1
    
    # 保存
    prs.save(output_path)
    print(f"\n✅ 专业 PPT 已保存：{output_path}")
    print(f"📊 共 {total_pages} 页幻灯片")
    print(f"🎨 使用主题：{theme['name']}")
    print(f"📐 遵循设计原则：对齐 · 对比 · 重复 · 亲密性")
    
    return output_path

if __name__ == '__main__':
    output_path = '/Users/yangbowen/Desktop/OpenClaw_企业应用方案_专业版.pptx'
    create_professional_ppt(output_path)
