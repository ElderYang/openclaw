#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenClaw 企业应用方案 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 创建 PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
COLORS = {
    'primary': RGBColor(25, 118, 210),      # 蓝色主色
    'secondary': RGBColor(0, 188, 212),     # 青色
    'accent': RGBColor(255, 152, 0),        # 橙色强调
    'success': RGBColor(76, 175, 80),       # 绿色
    'warning': RGBColor(255, 193, 7),       # 黄色
    'danger': RGBColor(244, 67, 54),        # 红色
    'dark': RGBColor(33, 33, 33),           # 深灰
    'light': RGBColor(245, 245, 245),       # 浅灰
    'white': RGBColor(255, 255, 255),       # 白色
}

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_points, layout_index=1):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    slide.shapes.title.text = title
    
    # 添加内容
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point['text']
        p.level = point.get('level', 0)
        p.font.size = Pt(18)
        if point.get('bold'):
            p.font.bold = True
    
    return slide

def add_comparison_slide(prs, title, items):
    """添加对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    
    # 创建表格
    rows = len(items) + 1
    cols = len(items[0]['columns']) if items else 0
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.333)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置列宽
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width
    
    # 填充表头
    if items:
        for i, col in enumerate(items[0]['columns']):
            cell = table.cell(0, i)
            cell.text = col
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['primary']
            cell.text_frame.paragraphs[0].font.color.rgb = COLORS['white']
            cell.text_frame.paragraphs[0].font.bold = True
    
    # 填充内容
    for row_idx, item in enumerate(items, 1):
        for col_idx, value in enumerate(item['values']):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light']
    
    return slide

def add_diagram_slide(prs, title, diagram_text):
    """添加图表页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    
    # 添加文本框显示图表
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.333)
    height = Inches(5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = diagram_text
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    
    return slide

# ==================== 创建 PPT ====================

print("创建 OpenClaw 企业应用方案 PPT...")

# 1. 标题页
print("1/15 创建标题页...")
slide = add_title_slide(
    prs,
    "OpenClaw 企业应用方案",
    f"技术架构 · 应用场景 · 部署方案\n中国移动终端公司\n{datetime.now().strftime('%Y年%m月')}"
)

# 2. 目录
print("2/15 创建目录页...")
slide = add_content_slide(
    prs,
    "目录",
    [
        {'text': 'OpenClaw 技术架构与原理', 'level': 0, 'bold': True},
        {'text': 'OpenClaw vs 对话式 AI vs 普通 Agent vs RPA', 'level': 0, 'bold': True},
        {'text': '中国移动终端公司应用场景', 'level': 0, 'bold': True},
        {'text': 'OpenClaw + RPA 协同方案', 'level': 0, 'bold': True},
        {'text': '为一线创造的价值', 'level': 0, 'bold': True},
        {'text': '内网部署硬件要求（70B+ 大模型）', 'level': 0, 'bold': True},
        {'text': '内网环境限制与风险', 'level': 0, 'bold': True},
        {'text': '总结与建议', 'level': 0, 'bold': True},
    ]
)

# 3. OpenClaw 技术架构
print("3/15 创建技术架构页...")
slide = add_diagram_slide(
    prs,
    "OpenClaw 技术架构",
    """
┌─────────────────────────────────────────────────────────────────┐
│                        用户交互层                                │
│         飞书/微信/Telegram/Discord/网页                          │
└─────────────────────────────────────────────────────────────────┘
                              ↓
┌─────────────────────────────────────────────────────────────────┐
│                      Gateway 网关层                              │
│    消息路由 · 会话管理 · 安全认证 · 限流控制                      │
└─────────────────────────────────────────────────────────────────┘
                              ↓
┌─────────────────────────────────────────────────────────────────┐
│                     Agent 核心层 (Pi-Agent)                      │
│   意图识别 · 任务规划 · 工具调用 · 记忆检索 · 结果生成            │
└─────────────────────────────────────────────────────────────────┘
                              ↓
┌─────────────────────────────────────────────────────────────────┐
│                      技能工具层                                  │
│   文件操作 · 浏览器控制 · API 调用 · 代码执行 · 数据库查询          │
└─────────────────────────────────────────────────────────────────┘
                              ↓
┌─────────────────────────────────────────────────────────────────┐
│                      记忆系统                                    │
│   短期记忆 (Session) · 长期记忆 (MEMORY.md) · 向量检索           │
└─────────────────────────────────────────────────────────────────┘
    """
)

# 4. OpenClaw 技术原理
print("4/15 创建技术原理页...")
slide = add_content_slide(
    prs,
    "OpenClaw 技术原理",
    [
        {'text': '核心设计理念', 'level': 0, 'bold': True},
        {'text': '• 本地优先（Local-First）：所有数据存储在本地，保护隐私', 'level': 1},
        {'text': '• 极简内核 + 弹性扩展：核心代码精简，通过 Skills 扩展能力', 'level': 1},
        {'text': '• 多模型兼容：支持百炼、Gemini、GPT、Claude 等主流模型', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Agent 循环（Agent Loop）', 'level': 0, 'bold': True},
        {'text': '1. 感知：接收用户消息，理解上下文', 'level': 1},
        {'text': '2. 思考：LLM 推理，生成行动计划', 'level': 1},
        {'text': '3. 行动：调用工具/Skills 执行任务', 'level': 1},
        {'text': '4. 反馈：记录结果，更新记忆', 'level': 1},
        {'text': '5. 学习：从反馈中优化策略', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '关键技术', 'level': 0, 'bold': True},
        {'text': '• 两层持久化：Session 状态 + 长期记忆', 'level': 1},
        {'text': '• 树状会话结构：支持多轮对话和上下文切换', 'level': 1},
        {'text': '• 压缩前落盘：长上下文自动压缩存储', 'level': 1},
    ]
)

# 5. 四类技术对比
print("5/15 创建对比页...")
slide = add_comparison_slide(
    prs,
    "OpenClaw vs 对话式 AI vs 普通 Agent vs RPA",
    [
        {'columns': ['对比维度', '对话式 AI', '普通 Agent', 'RPA', 'OpenClaw'], 'values': []},
        {'columns': [], 'values': ['核心能力', '文本对话', '简单任务', '规则自动化', '系统级任务执行']},
        {'columns': [], 'values': ['上下文理解', '短期对话', '有限', '无', '长期记忆']},
        {'columns': [], 'values': ['工具调用', '❌', '部分', '固定脚本', '✅ 丰富 Skills']},
        {'columns': [], 'values': ['学习能力', '❌', '弱', '❌', '✅ 自我改进']},
        {'columns': [], 'values': ['主动性', '被动响应', '被动', '定时触发', '主动代理']},
        {'columns': [], 'values': ['部署方式', '云端', '云端/本地', '本地', '本地优先']},
        {'columns': [], 'values': ['数据隐私', '低', '中', '高', '高']},
        {'columns': [], 'values': ['维护成本', '低', '中', '高（脚本维护）', '中']},
        {'columns': [], 'values': ['适用场景', '客服问答', '简单查询', '重复流程', '复杂任务']},
    ]
)

# 6. 中国移动终端公司应用场景
print("6/15 创建应用场景页...")
slide = add_content_slide(
    prs,
    "中国移动终端公司应用场景",
    [
        {'text': '市场部', 'level': 0, 'bold': True},
        {'text': '• 竞品分析：自动抓取竞品价格、配置、营销信息', 'level': 1},
        {'text': '• 舆情监控：实时监测社交媒体品牌提及', 'level': 1},
        {'text': '• 报告生成：自动生成销售日报、周报', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '产品部', 'level': 0, 'bold': True},
        {'text': '• 需求分析：自动整理用户反馈和评论', 'level': 1},
        {'text': '• 竞品对比：自动生成产品对比表', 'level': 1},
        {'text': '• 文档管理：自动整理产品文档和规格书', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '客服部', 'level': 0, 'bold': True},
        {'text': '• 智能问答：自动回答常见问题', 'level': 1},
        {'text': '• 工单处理：自动分类和分配工单', 'level': 1},
        {'text': '• 满意度分析：自动分析客户反馈', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '财务部', 'level': 0, 'bold': True},
        {'text': '• 发票处理：自动识别和录入发票', 'level': 1},
        {'text': '• 报表生成：自动生成财务报表', 'level': 1},
        {'text': '• 对账核对：自动对账和异常检测', 'level': 1},
    ]
)

# 7. 更多应用场景
print("7/15 创建更多应用场景页...")
slide = add_content_slide(
    prs,
    "中国移动终端公司应用场景（续）",
    [
        {'text': '供应链管理部', 'level': 0, 'bold': True},
        {'text': '• 库存监控：实时监控库存水平，自动预警', 'level': 1},
        {'text': '• 供应商管理：自动跟踪供应商交期和质量', 'level': 1},
        {'text': '• 采购分析：自动生成采购分析报告', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '人力资源部', 'level': 0, 'bold': True},
        {'text': '• 简历筛选：自动筛选和分类简历', 'level': 1},
        {'text': '• 培训管理：自动安排培训和跟踪进度', 'level': 1},
        {'text': '• 员工问答：自动回答 HR 政策问题', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'IT 运维部', 'level': 0, 'bold': True},
        {'text': '• 监控告警：自动监控系统状态，发送告警', 'level': 1},
        {'text': '• 日志分析：自动分析日志，发现异常', 'level': 1},
        {'text': '• 自动化运维：自动执行常见运维任务', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '研发部', 'level': 0, 'bold': True},
        {'text': '• 代码审查：自动审查代码规范', 'level': 1},
        {'text': '• 文档生成：自动生成技术文档', 'level': 1},
        {'text': '• Bug 分析：自动分析和分类 Bug', 'level': 1},
    ]
)

# 8. OpenClaw + RPA 协同方案
print("8/15 创建协同方案页...")
slide = add_diagram_slide(
    prs,
    "OpenClaw + RPA 协同方案",
    """
┌─────────────────────────────────────────────────────────────────┐
│                    发票处理混合架构示例                          │
└─────────────────────────────────────────────────────────────────┘

1️⃣ RPA 处理（确定性任务）
   └─ 从固定位置下载发票 PDF
       ↓
2️⃣ OpenClaw 处理（语义理解）
   └─ 理解发票内容，提取关键字段
       ↓
3️⃣ RPA 处理（确定性任务）
   └─ 将数据写入 ERP 系统
       ↓
4️⃣ OpenClaw 处理（异常处理）
   └─ 处理异常情况，生成报告

┌─────────────────────────────────────────────────────────────────┐
│                        职责分工                                  │
├─────────────────────────────────────────────────────────────────┤
│  RPA 擅长：            │  OpenClaw 擅长：                        │
│  • 固定流程操作         │  • 语义理解                             │
│  • 系统间数据搬运       │  • 决策判断                             │
│  • 批量数据处理         │  • 异常处理                             │
│  • 定时任务执行         │  • 自然语言交互                         │
└─────────────────────────────────────────────────────────────────┘
    """
)

# 9. 为一线创造的价值
print("9/15 创建价值页...")
slide = add_content_slide(
    prs,
    "为一线创造的价值",
    [
        {'text': '效率提升', 'level': 0, 'bold': True},
        {'text': '• 自动化重复工作：节省 50%+ 手工操作时间', 'level': 1},
        {'text': '• 7x24 小时工作：无需休息，随时响应', 'level': 1},
        {'text': '• 批量处理：一次性处理大量任务', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '质量提升', 'level': 0, 'bold': True},
        {'text': '• 减少人为错误：自动化流程零差错', 'level': 1},
        {'text': '• 标准化输出：统一格式和质量', 'level': 1},
        {'text': '• 可追溯：完整记录所有操作', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '成本降低', 'level': 0, 'bold': True},
        {'text': '• 人力成本：减少重复性工作岗位', 'level': 1},
        {'text': '• 培训成本：新员工快速上手', 'level': 1},
        {'text': '• 错误成本：减少返工和损失', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '员工体验', 'level': 0, 'bold': True},
        {'text': '• 解放创造力：从重复工作中解放', 'level': 1},
        {'text': '• 技能提升：学习 AI 工具使用', 'level': 1},
        {'text': '• 工作满意度：专注更有价值的工作', 'level': 1},
    ]
)

# 10. 量化价值
print("10/15 创建量化价值页...")
slide = add_comparison_slide(
    prs,
    "量化价值评估",
    [
        {'columns': ['场景', '当前耗时', 'OpenClaw 后', '节省时间', '年化价值'], 'values': []},
        {'columns': [], 'values': ['销售日报', '2 小时/天', '10 分钟/天', '95%', '约 5 万元']},
        {'columns': [], 'values': ['竞品分析', '8 小时/周', '1 小时/周', '87%', '约 15 万元']},
        {'columns': [], 'values': ['发票录入', '4 小时/天', '30 分钟/天', '85%', '约 8 万元']},
        {'columns': [], 'values': ['客服问答', '6 小时/天', '1 小时/天', '83%', '约 12 万元']},
        {'columns': [], 'values': ['库存监控', '2 小时/天', '10 分钟/天', '92%', '约 4 万元']},
        {'columns': [], 'values': ['合计', '-', '-', '-', '约 44 万元/年']},
    ]
)

# 11. 内网部署硬件要求
print("11/15 创建硬件要求页...")
slide = add_content_slide(
    prs,
    "内网部署硬件要求（70B+ 大模型）",
    [
        {'text': 'GPU 配置（核心）', 'level': 0, 'bold': True},
        {'text': '• 最低配置：4x NVIDIA A100 80GB 或 8x NVIDIA A800 80GB', 'level': 1},
        {'text': '• 推荐配置：8x NVIDIA H100 80GB 或 8x NVIDIA H800 80GB', 'level': 1},
        {'text': '• 替代方案：16x NVIDIA RTX 4090 24GB（成本更低）', 'level': 1},
        {'text': '• 显存要求：70B 模型至少需要 140GB 显存（FP16 精度）', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'CPU 配置', 'level': 0, 'bold': True},
        {'text': '• 最低：2x AMD EPYC 7763 或 Intel Xeon Gold 6348', 'level': 1},
        {'text': '• 推荐：4x AMD EPYC 9654 或 Intel Xeon Platinum 8480+', 'level': 1},
        {'text': '• 核心数：至少 64 核心，推荐 128+ 核心', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '内存配置', 'level': 0, 'bold': True},
        {'text': '• 最低：512GB DDR4/DDR5', 'level': 1},
        {'text': '• 推荐：1TB-2TB DDR5', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '存储配置', 'level': 0, 'bold': True},
        {'text': '• 系统盘：2TB NVMe SSD', 'level': 1},
        {'text': '• 数据盘：10TB+ NVMe SSD（模型和缓存）', 'level': 1},
        {'text': '• 备份盘：20TB+ HDD（冷备份）', 'level': 1},
    ]
)

# 12. 硬件成本估算
print("12/15 创建硬件成本页...")
slide = add_comparison_slide(
    prs,
    "硬件成本估算（70B 大模型）",
    [
        {'columns': ['配置方案', 'GPU', 'CPU', '内存', '存储', '总成本'], 'values': []},
        {'columns': [], 'values': ['入门级', '4x A800 80GB', '2x EPYC 7763', '512GB', '12TB', '约 80 万元']},
        {'columns': [], 'values': ['推荐级', '8x H800 80GB', '4x EPYC 9654', '1TB', '20TB', '约 200 万元']},
        {'columns': [], 'values': ['高性能', '8x H100 80GB', '4x Xeon Platinum', '2TB', '30TB', '约 350 万元']},
        {'columns': [], 'values': ['替代方案', '16x RTX 4090', '2x EPYC 7763', '512GB', '12TB', '约 50 万元']},
        {'columns': [], 'values': ['', '', '', '', '', '']},
        {'columns': [], 'values': ['年度运营成本', '电费：约 10 万元/年', '维护：约 5 万元/年', '', '', '']},
    ]
)

# 13. 内网环境限制
print("13/15 创建限制页...")
slide = add_content_slide(
    prs,
    "内网环境限制与风险",
    [
        {'text': '网络限制', 'level': 0, 'bold': True},
        {'text': '• 无法访问外网 API：需要使用本地模型或代理', 'level': 1},
        {'text': '• 技能更新困难：需要手动导入新 Skills', 'level': 1},
        {'text': '• 模型更新滞后：无法实时获取最新模型', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '安全风险', 'level': 0, 'bold': True},
        {'text': '• 数据泄露：Agent 可能意外发送敏感数据', 'level': 1},
        {'text': '• 权限滥用：Agent 可能执行未授权操作', 'level': 1},
        {'text': '• 注入攻击：恶意输入可能控制 Agent 行为', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '合规风险', 'level': 0, 'bold': True},
        {'text': '• 数据合规：需要符合数据安全法要求', 'level': 1},
        {'text': '• 审计要求：需要完整操作日志', 'level': 1},
        {'text': '• 责任界定：AI 决策失误的责任归属', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '运维风险', 'level': 0, 'bold': True},
        {'text': '• 系统依赖：过度依赖 AI 可能影响业务连续性', 'level': 1},
        {'text': '• 技能缺口：缺乏 AI 运维人才', 'level': 1},
        {'text': '• 模型漂移：模型性能可能随时间下降', 'level': 1},
    ]
)

# 14. 风险缓解方案
print("14/15 创建风险缓解页...")
slide = add_content_slide(
    prs,
    "风险缓解方案",
    [
        {'text': '安全加固', 'level': 0, 'bold': True},
        {'text': '• 沙箱隔离：所有操作在容器内执行', 'level': 1},
        {'text': '• 权限控制：最小权限原则，操作审批', 'level': 1},
        {'text': '• 敏感词过滤：自动检测和拦截敏感信息', 'level': 1},
        {'text': '• 操作审计：完整记录所有操作日志', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '合规保障', 'level': 0, 'bold': True},
        {'text': '• 数据脱敏：敏感数据自动脱敏处理', 'level': 1},
        {'text': '• 本地部署：所有数据存储在本地', 'level': 1},
        {'text': '• 人工审核：关键操作需要人工确认', 'level': 1},
        {'text': '• 定期审计：定期安全审计和合规检查', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '运维保障', 'level': 0, 'bold': True},
        {'text': '• 冗余部署：关键服务双机热备', 'level': 1},
        {'text': '• 监控告警：实时监控和自动告警', 'level': 1},
        {'text': '• 应急预案：制定详细的应急预案', 'level': 1},
        {'text': '• 培训计划：定期培训和技能提升', 'level': 1},
    ]
)

# 15. 总结与建议
print("15/15 创建总结页...")
slide = add_content_slide(
    prs,
    "总结与建议",
    [
        {'text': '核心优势', 'level': 0, 'bold': True},
        {'text': '✅ 本地优先：数据完全可控，符合内网要求', 'level': 1},
        {'text': '✅ 系统级执行：超越对话，实现真实任务自动化', 'level': 1},
        {'text': '✅ 丰富生态：100+ Skills，覆盖常见场景', 'level': 1},
        {'text': '✅ 灵活扩展：支持自定义 Skills 和工具', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '实施建议', 'level': 0, 'bold': True},
        {'text': '1. 试点先行：选择 1-2 个部门试点（建议市场部/客服部）', 'level': 1},
        {'text': '2. 场景优先：从高频、重复、规则明确的场景入手', 'level': 1},
        {'text': '3. 人机协同：初期以辅助为主，逐步过渡到自动化', 'level': 1},
        {'text': '4. 安全先行：先部署安全加固，再上线业务功能', 'level': 1},
        {'text': '5. 持续优化：定期评估效果，持续优化和改进', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '下一步行动', 'level': 0, 'bold': True},
        {'text': '📋 需求调研：深入各部门了解具体需求', 'level': 1},
        {'text': '📋 方案设计：制定详细的技术方案和实施计划', 'level': 1},
        {'text': '📋 环境准备：准备硬件环境和网络环境', 'level': 1},
        {'text': '📋 试点部署：选择试点场景进行部署', 'level': 1},
        {'text': '📋 效果评估：评估试点效果，决定是否推广', 'level': 1},
    ]
)

# 保存 PPT
output_path = '/Users/yangbowen/.openclaw/workspace/output/OpenClaw_企业应用方案_20260310.pptx'
prs.save(output_path)
print(f"\n✅ PPT 已保存：{output_path}")
print(f"📊 共 15 页幻灯片")
